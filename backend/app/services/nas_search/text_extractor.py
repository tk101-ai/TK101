"""PDF/DOCX/PPTX/HWP/HWPX/XLSX 텍스트 추출 + 단순 청크 분할.

PoC 단계라 tiktoken 없이 문자 기반 근사를 쓴다.
- chunk_size: 약 1500자(한국어 기준 약 500 token에 근접)
- chunk_overlap: 약 150자

v0.7.0 — 한글/엑셀 추가:
- .hwp (HWP5 OLE): olefile로 PrvText 스트림(미리보기 ~10KB)만 추출. 본문 BinData는 복잡도 대비 가치 낮아 스킵.
- .hwpx: zip+xml 표준 포맷이라 zipfile + xml.etree로 hp:t 텍스트 노드 추출.
- .xlsx: openpyxl read_only 모드로 셀 값만 join. 큰 시트 폭주 방지를 위해 시트당 추출 길이 상한.
"""
from __future__ import annotations

import logging
import zipfile
from pathlib import Path
from xml.etree import ElementTree as ET

logger = logging.getLogger(__name__)

# 문자 기반 청크 크기. token 환산은 모델별 편차가 커서 PoC에서는 단순화.
DEFAULT_CHUNK_CHARS = 1500
DEFAULT_OVERLAP_CHARS = 150

# 큰 .xlsx에서 셀이 백만 단위로 많을 수 있어 시트당 텍스트 길이를 제한.
# 임베딩 스텝에서 수십 MB 텍스트가 풀려 메모리/벡터 폭주하는 것을 방지.
XLSX_PER_SHEET_CHAR_LIMIT = 100_000

# HWPX 본문 XML 네임스페이스 — 한글의 표준 OWPML 스펙.
# 텍스트 노드는 hp:t. 안전을 위해 namespace 무시 매칭도 fallback으로 둔다.
HWPX_HP_NS = "http://www.hancom.co.kr/hwpml/2011/paragraph"


def extract_text(path: str) -> str:
    """확장자별 추출기 디스패치. 실패 시 빈 문자열 반환(상위 레이어가 skip 처리)."""
    ext = Path(path).suffix.lower()
    try:
        if ext == ".pdf":
            text = _extract_pdf(path)
        elif ext == ".docx":
            text = _extract_docx(path)
        elif ext == ".pptx":
            text = _extract_pptx(path)
        elif ext == ".hwp":
            text = _extract_hwp(path)
        elif ext == ".hwpx":
            text = _extract_hwpx(path)
        elif ext == ".xlsx":
            text = _extract_xlsx(path)
        else:
            logger.debug("지원하지 않는 확장자: %s", path)
            return ""
    except Exception as exc:  # noqa: BLE001 - 외부 라이브러리 다양한 예외 흡수
        logger.warning("텍스트 추출 실패: %s (%s)", path, exc)
        return ""
    # PostgreSQL UTF8 컬럼은 NULL byte(\x00)를 거부 — 일부 PDF 폰트 인코딩에서
    # "단어1\x00단어2" 형태로 섞여 들어와 INSERT 실패 → 파일 단위 인덱싱 실패.
    return text.replace("\x00", "") if text else ""


def _extract_pdf(path: str) -> str:
    from pdfminer.high_level import extract_text as pdf_extract

    return (pdf_extract(path) or "").strip()


def _extract_docx(path: str) -> str:
    from docx import Document

    doc = Document(path)
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    # 표 안의 텍스트도 인덱싱 대상에 포함.
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts).strip()


def _extract_pptx(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                parts.append(text)
    return "\n".join(parts).strip()


def _extract_hwp(path: str) -> str:
    """HWP5 (OLE 컨테이너) 본문 추출.

    PrvText 스트림에서 한글 워드프로세서가 저장하는 미리보기 텍스트를 읽는다.
    - PrvText는 UTF-16LE 인코딩, 약 10KB 상한 (한글 자체 스펙).
    - 본문 BinData/Section* 스트림은 압축+자체 레코드 포맷이라 검색 ROI 대비 복잡도 과대 → 스킵.
    - PrvText만으로도 검색용 키워드 인덱싱은 충분 (제목/요약/주요 본문 일부 포함됨).
    """
    import olefile

    if not olefile.isOleFile(path):
        # HWP3 (구버전 비OLE) 등은 미지원. 빈 문자열 반환해 파일명 청크만 사용.
        logger.debug("OLE 형식이 아닌 .hwp — PrvText 추출 불가: %s", path)
        return ""

    ole = olefile.OleFileIO(path)
    try:
        if not ole.exists("PrvText"):
            logger.debug(".hwp에 PrvText 스트림 없음: %s", path)
            return ""
        with ole.openstream("PrvText") as stream:
            raw = stream.read()
    finally:
        ole.close()

    if not raw:
        return ""

    # UTF-16LE 디코드 (한글 스펙). 깨진 바이트는 무시하여 부분 텍스트라도 살린다.
    try:
        text = raw.decode("utf-16-le", errors="ignore")
    except UnicodeDecodeError:
        return ""
    return text.strip()


def _extract_hwpx(path: str) -> str:
    """HWPX (zip + OWPML XML) 본문 추출.

    구조: Contents/section0.xml, section1.xml ... 안의 hp:t 텍스트 노드 누적.
    네임스페이스(hp)는 한글 OWPML 표준이지만 일부 변종은 미선언이거나 다를 수 있어
    namespace를 떼어낸 local-name 비교를 fallback으로 둔다.
    """
    parts: list[str] = []
    with zipfile.ZipFile(path) as zf:
        # Contents/section{N}.xml — 본문 섹션. header/settings 등은 메타라 제외.
        section_names = [
            n
            for n in zf.namelist()
            if n.startswith("Contents/section") and n.endswith(".xml")
        ]
        if not section_names:
            return ""

        for name in sorted(section_names):
            with zf.open(name) as fh:
                try:
                    tree = ET.parse(fh)
                except ET.ParseError as exc:
                    logger.warning("HWPX section 파싱 실패: %s/%s (%s)", path, name, exc)
                    continue
            root = tree.getroot()
            # 1) 명시적 hp 네임스페이스 매칭.
            tagged = root.findall(f".//{{{HWPX_HP_NS}}}t")
            if tagged:
                for el in tagged:
                    if el.text:
                        parts.append(el.text)
                continue
            # 2) fallback: 네임스페이스 모르고 local-name이 't'인 모든 요소.
            for el in root.iter():
                local = el.tag.rsplit("}", 1)[-1] if "}" in el.tag else el.tag
                if local == "t" and el.text:
                    parts.append(el.text)

    return "\n".join(parts).strip()


def _extract_xlsx(path: str) -> str:
    """XLSX 셀 값만 텍스트화.

    - read_only=True + data_only=True: 수식 결과만, 메모리 효율 모드.
    - 시트당 XLSX_PER_SHEET_CHAR_LIMIT 도달 시 조기 종료 → 백만 셀 시트 폭주 차단.
    - 빈 셀/None은 스킵, 숫자는 str 변환.
    """
    from openpyxl import load_workbook

    wb = load_workbook(path, read_only=True, data_only=True)
    try:
        all_parts: list[str] = []
        for sheet in wb.worksheets:
            sheet_chars = 0
            sheet_parts: list[str] = []
            stop = False
            for row in sheet.iter_rows(values_only=True):
                if stop:
                    break
                for value in row:
                    if value is None:
                        continue
                    cell_text = str(value).strip()
                    if not cell_text:
                        continue
                    sheet_parts.append(cell_text)
                    sheet_chars += len(cell_text) + 1  # +1 ≈ 구분자 비용
                    if sheet_chars >= XLSX_PER_SHEET_CHAR_LIMIT:
                        logger.debug(
                            "XLSX 시트 길이 상한 도달 — 조기 종료: %s/%s",
                            path,
                            sheet.title,
                        )
                        stop = True
                        break
            if sheet_parts:
                # 시트별 헤더로 시트명도 검색 가능하게 prepend.
                all_parts.append(f"[{sheet.title}]")
                all_parts.extend(sheet_parts)
        return "\n".join(all_parts).strip()
    finally:
        wb.close()


def chunk_text(
    text: str,
    *,
    chunk_chars: int = DEFAULT_CHUNK_CHARS,
    overlap_chars: int = DEFAULT_OVERLAP_CHARS,
) -> list[str]:
    """텍스트를 chunk_chars 길이의 윈도우로 슬라이드. overlap만큼 겹쳐서 잘라냄.

    공백 정규화: 연속된 공백/개행을 단일 공백으로 압축해 청크 효율 향상.
    """
    if not text:
        return []
    if overlap_chars >= chunk_chars:
        overlap_chars = max(0, chunk_chars // 5)

    normalized = " ".join(text.split())
    if not normalized:
        return []

    chunks: list[str] = []
    start = 0
    length = len(normalized)
    step = chunk_chars - overlap_chars
    while start < length:
        end = min(start + chunk_chars, length)
        piece = normalized[start:end].strip()
        if piece:
            chunks.append(piece)
        if end >= length:
            break
        start += step
    return chunks
