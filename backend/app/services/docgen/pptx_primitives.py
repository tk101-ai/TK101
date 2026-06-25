"""pptx 렌더 공유 프리미티브 — 빌더와 레이아웃 아카이브가 함께 쓴다.

저수준 도형/텍스트 추가, 공통 그리드(여백), 제목바·푸터·캡션 같은 재사용
구성요소를 모은다. pptx_builder(표지·목차·구분·본문·표·차트)와 pptx_layouts
(kpi·compare·timeline·steps·quote 아카이브)가 같은 시각 토큰을 공유하도록 한다.
"""
from __future__ import annotations

from app.services.docgen.theme import Theme

# 좌우 여백(인치). 모든 슬라이드 공통 그리드.
_MARGIN_IN = 0.7


def _rgb(color: tuple[int, int, int]):
    from pptx.dml.color import RGBColor

    return RGBColor(*color)


def _blank_layout(prs):
    """플레이스홀더가 가장 적은(빈) 레이아웃 — 기본/커스텀 템플릿 모두 안전."""
    best, best_count = None, 999
    for lay in prs.slide_layouts:
        count = len(lay.placeholders)
        if count < best_count:
            best, best_count = lay, count
    return best or prs.slide_layouts[0]


def _add_rect(slide, x, y, w, h, color, line_color=None):
    from pptx.enum.shapes import MSO_SHAPE

    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(color)
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = _rgb(line_color)
    shape.shadow.inherit = False
    return shape


def _add_text(slide, x, y, w, h, text, *, size, color, font, bold=False,
              align=None, anchor=None, spacing=None, wrap=True):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt

    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    if anchor is not None:
        tf.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE,
                              "bottom": MSO_ANCHOR.BOTTOM}[anchor]
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER,
                       "right": PP_ALIGN.RIGHT}[align]
    if spacing is not None:
        p.line_spacing = spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = _rgb(color)
    return box


def _kicker_text(label: str) -> str:
    """짧은 키커 라벨을 대문자+자간 느낌(글자 사이 공백)으로 변환해 에디토리얼 톤."""
    s = (label or "").strip().upper()
    # ASCII 약어만 자간을 벌린다(한글은 그대로 — 자모 분리 방지).
    if s and all(ord(c) < 128 for c in s):
        return " ".join(s)
    return s


def _slide_dims(prs):
    return prs.slide_width, prs.slide_height


def _add_footer(slide, prs, theme: Theme, page_no: int) -> None:
    from pptx.util import Inches, Pt

    sw, sh = _slide_dims(prs)
    m = Inches(_MARGIN_IN)
    # 옅은 구분선(헤어라인).
    _add_rect(slide, m, sh - Inches(0.62), sw - Inches(_MARGIN_IN * 2), Pt(0.75),
              theme.hairline)
    if theme.footer_text:
        _add_text(slide, m, sh - Inches(0.6), Inches(6), Inches(0.35),
                  theme.footer_text, size=9, color=theme.muted, font=theme.body_font,
                  align="left", anchor="middle")
    _add_text(slide, sw - Inches(1.6), sh - Inches(0.6), Inches(1.0), Inches(0.35),
              str(page_no), size=10, color=theme.accent, font=theme.heading_font,
              bold=True, align="right", anchor="middle")


def _add_title_bar(slide, prs, theme: Theme, heading: str, *, kicker: str = "",
                   continued: bool = False) -> None:
    """본문 슬라이드 상단: 섹션 키커 + 제목 + 액센트 밑줄."""
    from pptx.util import Inches, Pt

    sw, _ = _slide_dims(prs)
    m = Inches(_MARGIN_IN)
    width = sw - Inches(_MARGIN_IN * 2)
    if kicker:
        _add_text(slide, m + Inches(0.02), Inches(0.42), width, Inches(0.3),
                  _kicker_text(kicker), size=11, color=theme.accent,
                  font=theme.heading_font, bold=True, align="left", anchor="middle")
    text = heading or " "
    if continued:
        text += "  (계속)"
    top = Inches(0.72) if kicker else Inches(0.5)
    _add_text(slide, m, top, width, Inches(0.7), text, size=25, color=theme.primary,
              font=theme.heading_font, bold=True, align="left", anchor="middle")
    _add_rect(slide, m + Inches(0.02), Inches(1.42), Inches(1.7), Pt(3.5), theme.accent)


def _add_caption(slide, prs, theme: Theme, text: str) -> None:
    """제목바 아래 짧은 리드 문장(표/차트/레이아웃 위 캡션)."""
    from pptx.util import Inches

    sw, _ = _slide_dims(prs)
    _add_text(slide, Inches(_MARGIN_IN), Inches(1.52), sw - Inches(_MARGIN_IN * 2),
              Inches(0.4), text, size=14, color=theme.muted, font=theme.body_font,
              align="left", anchor="middle")
