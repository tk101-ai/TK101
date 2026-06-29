"""구조화 섹션 → 디자인된 .pptx 바이트 (고품질 렌더).

16:9 + 브랜드 테마(theme.py)로 다음을 구성한다:
- 레이어드 표지(짙은 primary 톤블록 + 액센트 키커/룰)
- 목차(번호·헤어라인)
- 섹션 구분 슬라이드(거대 번호 + 제목, 덱에 리듬 부여)
- 본문 슬라이드(섹션 키커 + 제목 + 액센트 밑줄, 색 글리프 불릿)
- 정돈된 표(그리드 제거 + 헤더 룰 + 줄무늬)
- 브랜드 색 파생 차트(시리즈 팔레트 + 옅은 그리드라인)
- 푸터(헤어라인 + 브랜드 + 페이지)

body 의 불릿/표/차트는 markdown_blocks.parse_blocks 로 해석한다. 회사 .pptx 템플릿
(settings.docgen_pptx_template)이 있으면 그 마스터를 베이스로 쓴다.

자동 생성물이라 **절대 렌더가 통째로 깨지면 안 된다** — 장식/차트 스타일은 모두
try/except 로 감싸 실패해도 슬라이드는 남긴다.
"""
from __future__ import annotations

import io
import logging
from datetime import date

from app.services.docgen.markdown_blocks import Block, parse_blocks
from app.services.docgen.pptx_layouts import render_layout_slide
from app.services.docgen.pptx_primitives import (
    _MARGIN_IN,
    _add_caption,
    _add_footer,
    _add_rect,
    _add_text,
    _add_title_bar,
    _blank_layout,
    _kicker_text,
    _rgb,
    _slide_dims,
)
from app.services.docgen.theme import Theme, get_theme

logger = logging.getLogger(__name__)

# 슬라이드 한 장에 들어갈 본문 텍스트 줄 수 상한(넘으면 다음 장으로 분할).
_MAX_TEXT_LINES = 9
# 표 한 장 최대 행 수(헤더 포함). 넘으면 분할.
_MAX_TABLE_ROWS = 12
# 레벨별 불릿 글리프(level 0 은 색 사각형, 이하 대시/점).
_BULLET_GLYPHS = ("■", "–", "·", "·", "·")
# 섹션이 2개 이상이면 섹션 구분 슬라이드를 넣어 덱에 리듬을 준다.
_SECTION_DIVIDERS = True
# 그리드라인 없는 깔끔한 표 스타일(Office "No Style, No Grid").
_TABLE_STYLE_NO_GRID = "{2D5ABB26-0587-4C30-8999-92F81FD0307C}"


def _cover_slide(prs, theme: Theme, title: str, subtitle: str, kicker: str) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    sw, sh = _slide_dims(prs)
    # 풀블리드 배경 + 우측 톤블록(레이어 깊이) + 액센트 밴드.
    _add_rect(slide, 0, 0, sw, sh, theme.primary)
    _add_rect(slide, int(sw * 0.66), 0, sw - int(sw * 0.66), sh, theme.primary_deep)
    _add_rect(slide, 0, sh - Inches(0.16), sw, Inches(0.16), theme.accent)
    # 제목 위 짧고 굵은 액센트 바.
    _add_rect(slide, Inches(0.72), Inches(2.55), Inches(0.9), Pt(6), theme.accent)
    # 로고(있으면 좌상단).
    if theme.logo_path:
        try:
            slide.shapes.add_picture(theme.logo_path, Inches(0.72), Inches(0.6),
                                     height=Inches(0.7))
        except Exception:  # noqa: BLE001 - 로고 깨져도 표지는 나와야 함
            pass
    _cover_sub = (0xD7, 0xDE, 0xEA)  # 진한 배경 위 흐린 흰색.
    if kicker:
        _add_text(slide, Inches(0.72), Inches(2.0), sw - Inches(1.4), Inches(0.4),
                  _kicker_text(kicker), size=14, color=theme.accent_soft,
                  font=theme.heading_font, bold=True, align="left", anchor="middle")
    _add_text(slide, Inches(0.7), Inches(2.85), sw - Inches(1.4), Inches(2.0),
              title or "문서", size=42, color=theme.white, font=theme.heading_font,
              bold=True, align="left", anchor="top", spacing=1.05)
    # 제목과 부제 사이 가는 구분 룰.
    _add_rect(slide, Inches(0.74), Inches(4.95), Inches(3.4), Pt(1.0), _cover_sub)
    if subtitle:
        _add_text(slide, Inches(0.72), Inches(5.15), sw - Inches(1.4), Inches(0.5),
                  subtitle, size=15, color=_cover_sub, font=theme.body_font,
                  align="left")
    if theme.footer_text:
        _add_text(slide, Inches(0.72), sh - Inches(0.95), sw - Inches(1.4),
                  Inches(0.5), theme.footer_text, size=12, color=_cover_sub,
                  font=theme.body_font, align="left", anchor="middle")


def _agenda_slide(prs, theme: Theme, headings: list[str], page_no: int) -> None:
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    _add_title_bar(slide, prs, theme, "목차", kicker="CONTENTS")
    sw, _ = _slide_dims(prs)
    m = Inches(_MARGIN_IN)
    top = Inches(1.85)
    row_h = min(0.66, 4.6 / max(len(headings), 1))
    for idx, h in enumerate(headings, 1):
        y = top + Inches(row_h) * (idx - 1)
        # 큰 액센트 번호.
        _add_text(slide, m, y, Inches(0.8), Inches(row_h), f"{idx:02d}",
                  size=22, color=theme.accent, font=theme.heading_font, bold=True,
                  align="left", anchor="middle")
        _add_text(slide, m + Inches(0.95), y, sw - Inches(_MARGIN_IN * 2) - Inches(0.95),
                  Inches(row_h), h, size=16, color=theme.text, font=theme.body_font,
                  anchor="middle")
        # 항목 사이 헤어라인.
        if idx < len(headings):
            _add_rect(slide, m, y + Inches(row_h) - Pt(0.5),
                      sw - Inches(_MARGIN_IN * 2), Pt(0.5), theme.hairline)
    _add_footer(slide, prs, theme, page_no)


def _section_divider_slide(prs, theme: Theme, index: int, total: int, heading: str,
                           page_no: int) -> None:
    """섹션 시작 구분 슬라이드 — 거대 번호 + 제목으로 덱에 리듬을 준다."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    sw, sh = _slide_dims(prs)
    m = Inches(_MARGIN_IN)
    # 중성 배경 + 좌측 액센트 룰.
    _add_rect(slide, 0, 0, sw, sh, theme.surface)
    _add_rect(slide, 0, 0, Inches(0.16), sh, theme.accent)
    # 우측에 거대한 옅은 번호(워터마크 느낌).
    _add_text(slide, sw - Inches(5.0), Inches(0.2), Inches(4.6), sh - Inches(0.4),
              f"{index:02d}", size=230, color=theme.accent_soft,
              font=theme.heading_font, bold=True, align="right", anchor="middle",
              wrap=False)
    # 키커 + 제목(좌측, 수직 중앙).
    _add_text(slide, m + Inches(0.1), Inches(2.7), Inches(7.5), Inches(0.4),
              _kicker_text(f"SECTION {index:02d} / {total:02d}"), size=13,
              color=theme.accent, font=theme.heading_font, bold=True, align="left")
    _add_rect(slide, m + Inches(0.12), Inches(3.2), Inches(0.9), Pt(5), theme.accent)
    _add_text(slide, m + Inches(0.1), Inches(3.4), Inches(8.0), Inches(1.6),
              heading or " ", size=34, color=theme.primary, font=theme.heading_font,
              bold=True, align="left", anchor="top", spacing=1.05)
    _add_footer(slide, prs, theme, page_no)


def _text_lines_from_blocks(blocks: list[Block]) -> list[tuple[int, str, bool]]:
    """문단/불릿 블록을 (level, text, is_bullet) 줄 목록으로 평탄화."""
    lines: list[tuple[int, str, bool]] = []
    for b in blocks:
        if b.kind == "paragraph":
            for ln in b.lines:
                if ln.strip():
                    lines.append((0, ln.strip(), False))
        elif b.kind == "bullets":
            for level, text in b.items:
                if text.strip():
                    lines.append((level, text.strip(), True))
    return lines


def _add_text_slide(prs, theme: Theme, heading: str, lines, *, kicker: str,
                    continued: bool, page_no: int) -> None:
    from pptx.enum.text import MSO_ANCHOR
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    _add_title_bar(slide, prs, theme, heading, kicker=kicker, continued=continued)
    sw, _ = _slide_dims(prs)
    box = slide.shapes.add_textbox(Inches(_MARGIN_IN), Inches(1.7),
                                   sw - Inches(_MARGIN_IN * 2), Inches(4.7))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for idx, (level, text, is_bullet) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        p.line_spacing = 1.2
        if is_bullet:
            glyph = _BULLET_GLYPHS[min(level, len(_BULLET_GLYPHS) - 1)]
            # 글리프와 본문을 분리한 런 — level 0 사각형은 액센트색으로 강조.
            g = p.add_run()
            g.text = ("    " * level) + glyph + "  "
            g.font.name = theme.heading_font
            g.font.size = Pt(13 if level == 0 else 12)
            g.font.color.rgb = _rgb(theme.accent if level == 0 else theme.muted)
            run = p.add_run()
            run.text = text
            run.font.size = Pt(17 if level == 0 else 15)
            run.font.bold = level == 0
            run.font.color.rgb = _rgb(theme.text if level == 0 else theme.muted)
            run.font.name = theme.body_font
        else:
            run = p.add_run()
            run.text = text
            run.font.size = Pt(16)
            run.font.color.rgb = _rgb(theme.text)
            run.font.name = theme.body_font
    _add_footer(slide, prs, theme, page_no)


def _style_table_no_grid(table) -> None:
    """표 기본 그리드/밴딩 제거(Swiss 느낌). 실패해도 표는 남긴다."""
    try:
        from pptx.oxml.ns import qn

        table.first_row = True
        table.horz_banding = False
        tblPr = table._tbl.tblPr  # noqa: SLF001 - python-pptx 표 스타일 직접 설정.
        style_el = tblPr.find(qn("a:tableStyleId"))
        if style_el is None:
            style_el = tblPr.makeelement(qn("a:tableStyleId"), {})
            tblPr.append(style_el)
        style_el.text = _TABLE_STYLE_NO_GRID
    except Exception:  # noqa: BLE001 - 스타일 XML 실패가 표를 막아선 안 됨
        logger.debug("표 그리드 제거 실패 — 기본 스타일 유지", exc_info=True)


def _add_table_slide(prs, theme: Theme, heading: str, rows: list[list[str]], *,
                     kicker: str, continued: bool, page_no: int,
                     intro: str = "") -> None:
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    if not rows:
        return
    slide = prs.slides.add_slide(_blank_layout(prs))
    _add_title_bar(slide, prs, theme, heading, kicker=kicker, continued=continued)
    if intro:
        _add_caption(slide, prs, theme, intro)
    sw, _ = _slide_dims(prs)
    ncol = max(len(r) for r in rows)
    nrow = len(rows)
    top_in = 2.05 if intro else 1.75
    left, top = Inches(_MARGIN_IN), Inches(top_in)
    width = sw - Inches(_MARGIN_IN * 2)
    height = Inches(min(0.52 * nrow, 4.7 - (0.3 if intro else 0)))
    table = slide.shapes.add_table(nrow, ncol, left, top, width, height).table
    _style_table_no_grid(table)
    table.rows[0].height = Inches(0.5)
    for i, r in enumerate(rows):
        is_header = i == 0
        for j in range(ncol):
            cell = table.cell(i, j)
            cell.text = r[j] if j < len(r) else ""
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Inches(0.12)
            cell.margin_right = Inches(0.1)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            cell.fill.solid()
            if is_header:
                cell.fill.fore_color.rgb = _rgb(theme.primary)
            elif i % 2 == 1:
                cell.fill.fore_color.rgb = _rgb(theme.surface)
            else:
                cell.fill.fore_color.rgb = _rgb(theme.white)
            for para in cell.text_frame.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                for run in para.runs:
                    run.font.size = Pt(12)
                    run.font.name = theme.body_font
                    run.font.bold = is_header or j == 0
                    run.font.color.rgb = _rgb(
                        theme.white if is_header else theme.text
                    )
    # 헤더 하단 액센트 룰(표를 또렷하게).
    _add_rect(slide, left, top + Inches(0.5) - Pt(2), width, Pt(2), theme.accent)
    _add_footer(slide, prs, theme, page_no)


def _chart_type(kind: str):
    from pptx.enum.chart import XL_CHART_TYPE

    return {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }.get(kind, XL_CHART_TYPE.COLUMN_CLUSTERED)


def _style_chart(chart, theme: Theme, ctype: str, nseries: int) -> None:
    """차트에 브랜드 색·폰트·옅은 그리드라인 적용. 실패해도 차트는 남긴다."""
    from pptx.enum.chart import XL_LABEL_POSITION, XL_LEGEND_POSITION
    from pptx.util import Pt

    palette = theme.series_palette
    try:
        chart.font.size = Pt(11)
        chart.font.name = theme.body_font
    except Exception:  # noqa: BLE001
        pass
    # 범례.
    try:
        if ctype == "pie" or nseries > 1:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = False
    except Exception:  # noqa: BLE001
        pass
    # 시리즈/포인트 색.
    try:
        plot = chart.plots[0]
        if ctype == "pie":
            series = plot.series[0]
            for i, pt in enumerate(series.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = _rgb(palette[i % len(palette)])
            plot.has_data_labels = True
            plot.data_labels.show_percentage = True
            plot.data_labels.show_value = False
            plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
            plot.data_labels.font.size = Pt(10)
            plot.data_labels.font.name = theme.body_font
        else:
            for i, series in enumerate(plot.series):
                col = palette[i % len(palette)]
                if ctype == "line":
                    series.format.line.color.rgb = _rgb(col)
                    series.format.line.width = Pt(2.25)
                else:
                    series.format.fill.solid()
                    series.format.fill.fore_color.rgb = _rgb(col)
    except Exception:  # noqa: BLE001
        logger.debug("차트 시리즈 색 적용 실패", exc_info=True)
    # 축 — 옅은 그리드라인 + 작은 라벨.
    if ctype != "pie":
        try:
            va = chart.value_axis
            va.has_major_gridlines = True
            va.major_gridlines.format.line.color.rgb = _rgb(theme.hairline)
            va.major_gridlines.format.line.width = Pt(0.5)
            va.tick_labels.font.size = Pt(10)
            va.tick_labels.font.name = theme.body_font
            ca = chart.category_axis
            ca.tick_labels.font.size = Pt(10)
            ca.tick_labels.font.name = theme.body_font
            ca.has_major_gridlines = False
        except Exception:  # noqa: BLE001
            pass


def _add_chart_slide(prs, theme: Theme, heading: str, data: dict, *, kicker: str,
                     page_no: int, intro: str = "") -> None:
    from pptx.chart.data import CategoryChartData
    from pptx.util import Inches

    slide = prs.slides.add_slide(_blank_layout(prs))
    title = data.get("title") or heading
    _add_title_bar(slide, prs, theme, title, kicker=kicker)
    if intro:
        _add_caption(slide, prs, theme, intro)
    sw, _ = _slide_dims(prs)
    top_in = 2.1 if intro else 1.75
    chart_h = 4.3 if intro else 4.6
    # 잘못된 LLM 차트 스펙(다중 시리즈 pie·비숫자 값 등)이 python-pptx 내부에서
    # 예외를 던져 덱 전체 렌더가 500 나는 걸 막는다. 실패 시 해당 차트만 건너뛰고
    # 플레이스홀더 텍스트로 대체 — 나머지 슬라이드는 정상 렌더.
    try:
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]
        for s in data["series"]:
            chart_data.add_series(s["name"], s["values"])
        graphic = slide.shapes.add_chart(
            _chart_type(data["type"]), Inches(_MARGIN_IN + 0.1), Inches(top_in),
            sw - Inches(_MARGIN_IN * 2 + 0.2), Inches(chart_h), chart_data,
        )
        chart = graphic.chart
        chart.has_title = False
        _style_chart(chart, theme, data["type"], len(data["series"]))
    except Exception:  # noqa: BLE001 - 차트 깨져도 덱은 나와야 함
        logger.warning("차트 렌더 실패 — 플레이스홀더로 대체: %s", title, exc_info=True)
        _add_text(slide, Inches(_MARGIN_IN + 0.1), Inches(top_in),
                  sw - Inches(_MARGIN_IN * 2 + 0.2), Inches(chart_h),
                  "[차트를 표시할 수 없습니다]", size=16, color=theme.muted,
                  font=theme.body_font, align="center", anchor="middle")
    _add_footer(slide, prs, theme, page_no)


def _render_section(prs, theme: Theme, heading: str, blocks: list[Block],
                    counter: list[int], *, kicker: str) -> None:
    """한 섹션 → 텍스트/표/차트 슬라이드. counter[0]=현재 페이지 번호(가변)."""
    pending_text: list[Block] = []
    first_text = True

    def flush_text(more_after: bool) -> None:
        nonlocal first_text
        lines = _text_lines_from_blocks(pending_text)
        pending_text.clear()
        if not lines:
            return
        # _MAX_TEXT_LINES 단위로 분할.
        for start in range(0, len(lines), _MAX_TEXT_LINES):
            chunk = lines[start:start + _MAX_TEXT_LINES]
            counter[0] += 1
            _add_text_slide(prs, theme, heading, chunk, kicker=kicker,
                            continued=not first_text, page_no=counter[0])
            first_text = False

    def take_intro() -> str:
        """표/차트 바로 위 짧은 리드 문단(≤2줄, 불릿 없음)은 별도 슬라이드 대신
        같은 슬라이드의 캡션으로 접어 빈 슬라이드를 막는다. 그 외엔 정상 flush."""
        lines = _text_lines_from_blocks(pending_text)
        if lines and len(lines) <= 2 and all(not is_b for _, _, is_b in lines):
            pending_text.clear()
            return " ".join(t for _, t, _ in lines)
        flush_text(True)
        return ""

    for b in blocks:
        if b.kind in ("paragraph", "bullets"):
            pending_text.append(b)
        elif b.kind == "table":
            intro = take_intro()
            rows = b.rows
            for start in range(0, len(rows), _MAX_TABLE_ROWS):
                chunk = rows[start:start + _MAX_TABLE_ROWS]
                # 분할 시 헤더 반복.
                if start > 0 and rows:
                    chunk = [rows[0]] + chunk
                counter[0] += 1
                _add_table_slide(prs, theme, heading, chunk, kicker=kicker,
                                 continued=start > 0, page_no=counter[0],
                                 intro=intro if start == 0 else "")
        elif b.kind == "chart":
            intro = take_intro()
            counter[0] += 1
            _add_chart_slide(prs, theme, heading, b.data, kicker=kicker,
                             page_no=counter[0], intro=intro)
        elif b.kind == "layout":
            intro = take_intro()
            counter[0] += 1
            render_layout_slide(prs, theme, heading, b.data, kicker=kicker,
                                page_no=counter[0], intro=intro)
    flush_text(False)
    # 본문이 완전히 비었던 섹션도 제목 슬라이드 한 장 남긴다.
    if first_text and not any(b.kind in ("table", "chart", "layout") for b in blocks):
        counter[0] += 1
        _add_text_slide(prs, theme, heading, [], kicker=kicker, continued=False,
                        page_no=counter[0])


def build_pptx(
    title: str, sections: list[dict], theme: Theme | None = None
) -> bytes:
    """title + sections([{heading, body}]) → 디자인된 .pptx 바이트.

    theme 가 주어지면 그 색·폰트로(디자인 프리셋), 없으면 회사 기본 테마.
    """
    from pptx import Presentation
    from pptx.util import Inches

    theme = theme or get_theme()
    prs = Presentation(theme.pptx_template) if theme.pptx_template else Presentation()
    # 16:9. (커스텀 템플릿이면 템플릿 크기를 존중해 덮어쓰지 않음.)
    if not theme.pptx_template:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

    today = date.today().strftime("%Y.%m.%d")
    cover_kicker = theme.footer_text or "DOCUMENT"
    _cover_slide(prs, theme, title, today, cover_kicker)

    headings = [(s.get("heading") or "").strip() for s in sections]
    headings = [h for h in headings if h]
    total = len(headings)
    counter = [1]  # 표지=1.
    if total >= 2:
        counter[0] += 1
        _agenda_slide(prs, theme, headings, counter[0])

    use_dividers = _SECTION_DIVIDERS and total >= 2
    sec_no = 0
    for s in sections:
        heading = (s.get("heading") or "").strip()
        body = (s.get("body") or "").strip()
        if heading:
            sec_no += 1
        kicker = f"SECTION {sec_no:02d}" if heading else ""
        if use_dividers and heading:
            counter[0] += 1
            _section_divider_slide(prs, theme, sec_no, total, heading, counter[0])
        _render_section(prs, theme, heading, parse_blocks(body), counter,
                        kicker=kicker)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
