"""고품질 슬라이드 레이아웃 아카이브 — 내용 성격에 맞는 디자인된 원형들.

markdown_blocks 가 ```layout {json}``` 펜스를 Block(kind="layout", data=...) 로
파싱하면, pptx_builder._render_section 이 이 모듈의 render_layout_slide 로 보낸다.
획일적인 "제목 + 불릿" 덱을 벗어나, 통계는 KPI 카드, 비교는 대조 패널, 일정은
타임라인, 절차는 번호 단계, 핵심 한 줄은 인용 콜아웃으로 그린다.

설계 원칙:
- 모든 색/폰트는 theme(브랜드 파생)에서만 — 하드코딩 없음.
- 자동 생성물이라 **절대 통째로 깨지면 안 된다** — 각 렌더는 try/except 로 감싸고
  실패 시 제목/푸터만 남은 슬라이드로 폴백한다(render_layout_slide 에서 처리).
- 좌표는 EMU 정수로 캐스팅(python-pptx 는 float 좌표를 거부할 수 있음).

다음 단계(별도 PR): 생성 LLM 이 내용에 맞는 layout 을 골라 이 펜스를 emit 하도록
프롬프트/스키마를 연결한다. 이 PR 은 렌더 아카이브(시각 퀄리티)만 확보한다.
"""
from __future__ import annotations

import logging

from app.services.docgen.pptx_primitives import (
    _MARGIN_IN,
    _add_caption,
    _add_footer,
    _add_rect,
    _add_text,
    _add_title_bar,
    _blank_layout,
    _kicker_text,
    _rgb,
    _slide_dims,
)
from app.services.docgen.theme import Theme

logger = logging.getLogger(__name__)

# 본문 콘텐츠 영역 하단 한계(인치). 푸터 위 여백 확보.
_AREA_BOTTOM_IN = 6.6


def render_layout_slide(prs, theme: Theme, heading: str, data: dict, *,
                        kicker: str, page_no: int, intro: str = "") -> None:
    """레이아웃 블록 1개 → 디자인된 슬라이드 1장. 알 수 없는 타입/실패는 안전 폴백."""
    data = data or {}
    kind = str(data.get("layout") or "").lower()

    # 인용은 제목바 없는 풀 슬라이드 콜아웃 — 별도 처리.
    if kind == "quote":
        _render_quote_slide(prs, theme, data, page_no)
        return

    slide = prs.slides.add_slide(_blank_layout(prs))
    title = data.get("title") or heading
    _add_title_bar(slide, prs, theme, title, kicker=kicker)
    if intro:
        _add_caption(slide, prs, theme, intro)
    top_in = 2.05 if intro else 1.75
    try:
        if kind == "kpi":
            _render_kpi(slide, prs, theme, data, top_in)
        elif kind == "compare":
            _render_compare(slide, prs, theme, data, top_in)
        elif kind == "timeline":
            _render_timeline(slide, prs, theme, data, top_in)
        elif kind == "steps":
            _render_steps(slide, prs, theme, data, top_in)
    except Exception:  # noqa: BLE001 - 레이아웃 렌더 실패가 덱을 막아선 안 됨
        logger.warning("레이아웃 렌더 실패(%s) — 제목만 남김", kind, exc_info=True)
    _add_footer(slide, prs, theme, page_no)


def _add_bullets_box(slide, x, y, w, h, theme: Theme, points: list[str]) -> None:
    """패널 내부 불릿 목록 — 액센트 사각 글리프 + 본문."""
    from pptx.util import Pt

    box = slide.shapes.add_textbox(int(x), int(y), int(w), int(h))
    tf = box.text_frame
    tf.word_wrap = True
    for idx, text in enumerate(points):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(7)
        p.line_spacing = 1.15
        g = p.add_run()
        g.text = "■  "
        g.font.name = theme.heading_font
        g.font.size = Pt(11)
        g.font.color.rgb = _rgb(theme.accent)
        r = p.add_run()
        r.text = text
        r.font.name = theme.body_font
        r.font.size = Pt(14)
        r.font.color.rgb = _rgb(theme.text)


def _render_kpi(slide, prs, theme: Theme, data: dict, top_in: float) -> None:
    """큰 숫자 KPI 카드 — 통계/지표 강조(1~4개 가로 배치)."""
    from pptx.util import Inches, Pt

    items = data.get("items") or []
    n = len(items)
    if n == 0:
        return
    sw, _ = _slide_dims(prs)
    left = int(Inches(_MARGIN_IN))
    content_w = int(sw - Inches(_MARGIN_IN * 2))
    gap = int(Inches(0.3))
    card_w = (content_w - gap * (n - 1)) // n
    card_h = int(Inches(2.7))
    area_top = int(Inches(top_in))
    area_bottom = int(Inches(_AREA_BOTTOM_IN))
    card_top = area_top + (area_bottom - area_top - card_h) // 2
    for idx, it in enumerate(items):
        x = left + (card_w + gap) * idx
        _add_rect(slide, x, card_top, card_w, card_h, theme.surface)
        _add_rect(slide, x, card_top, card_w, int(Pt(5)), theme.accent)  # 상단 액센트.
        _add_text(slide, x, card_top + Inches(0.45), card_w, Inches(1.15),
                  it.get("value", ""), size=40, color=theme.primary,
                  font=theme.heading_font, bold=True, align="center", anchor="middle")
        _add_text(slide, x + int(Inches(0.1)), card_top + Inches(1.65),
                  card_w - int(Inches(0.2)), Inches(0.5), it.get("label", ""),
                  size=15, color=theme.text, font=theme.body_font, bold=True,
                  align="center", anchor="middle")
        caption = it.get("caption")
        if caption:
            _add_text(slide, x + int(Inches(0.1)), card_top + Inches(2.1),
                      card_w - int(Inches(0.2)), Inches(0.4), caption, size=11,
                      color=theme.muted, font=theme.body_font, align="center")


def _render_compare(slide, prs, theme: Theme, data: dict, top_in: float) -> None:
    """2~3단 대조 패널 — 기존/제안, 옵션 비교 등."""
    from pptx.util import Inches

    cols = data.get("columns") or []
    n = len(cols)
    if n == 0:
        return
    sw, _ = _slide_dims(prs)
    left = int(Inches(_MARGIN_IN))
    content_w = int(sw - Inches(_MARGIN_IN * 2))
    gap = int(Inches(0.35))
    col_w = (content_w - gap * (n - 1)) // n
    top = int(Inches(top_in))
    panel_h = int(Inches(_AREA_BOTTOM_IN)) - top
    head_h = int(Inches(0.62))
    for idx, col in enumerate(cols):
        x = left + (col_w + gap) * idx
        _add_rect(slide, x, top, col_w, panel_h, theme.surface)         # 패널 배경.
        _add_rect(slide, x, top, col_w, head_h, theme.primary)          # 헤더 밴드.
        _add_text(slide, x + int(Inches(0.1)), top, col_w - int(Inches(0.2)), head_h,
                  col.get("heading", ""), size=16, color=theme.white,
                  font=theme.heading_font, bold=True, align="center", anchor="middle")
        _add_bullets_box(slide, x + Inches(0.25), top + head_h + int(Inches(0.2)),
                         col_w - int(Inches(0.5)), panel_h - head_h - int(Inches(0.4)),
                         theme, col.get("points") or [])


def _render_timeline(slide, prs, theme: Theme, data: dict, top_in: float) -> None:
    """가로 타임라인/로드맵 — 단계별 마일스톤(2~5개)."""
    from pptx.util import Inches, Pt

    ms = data.get("milestones") or []
    n = len(ms)
    if n == 0:
        return
    sw, _ = _slide_dims(prs)
    left = int(Inches(_MARGIN_IN))
    content_w = int(sw - Inches(_MARGIN_IN * 2))
    slot_w = content_w // n
    base_y = int(Inches(top_in + 1.2))           # 노드 라인 y.
    # 노드들을 잇는 가로 룰.
    _add_rect(slide, left + slot_w // 2, base_y, content_w - slot_w, int(Pt(2.5)),
              theme.accent)
    node = int(Inches(0.22))
    for idx, m in enumerate(ms):
        cx = left + slot_w * idx + slot_w // 2
        # 노드(액센트 사각형).
        _add_rect(slide, cx - node // 2, base_y - node // 2, node, node, theme.accent)
        # 위: 라벨 + 제목.
        _add_text(slide, cx - slot_w // 2 + int(Inches(0.1)), int(Inches(top_in)),
                  slot_w - int(Inches(0.2)), Inches(0.35),
                  _kicker_text(m.get("label") or f"{idx + 1}단계"), size=11,
                  color=theme.accent, font=theme.heading_font, bold=True, align="center")
        _add_text(slide, cx - slot_w // 2 + int(Inches(0.1)),
                  int(Inches(top_in)) + int(Inches(0.4)), slot_w - int(Inches(0.2)),
                  Inches(0.6), m.get("title", ""), size=15, color=theme.primary,
                  font=theme.heading_font, bold=True, align="center", anchor="middle")
        # 아래: 상세.
        detail = m.get("detail")
        if detail:
            _add_text(slide, cx - slot_w // 2 + int(Inches(0.15)),
                      base_y + int(Inches(0.35)), slot_w - int(Inches(0.3)),
                      Inches(1.8), detail, size=12, color=theme.muted,
                      font=theme.body_font, align="center", spacing=1.1)


def _render_steps(slide, prs, theme: Theme, data: dict, top_in: float) -> None:
    """번호 단계 — 절차/방법론(2~6단계, 세로 배치)."""
    from pptx.util import Inches

    steps = data.get("steps") or []
    n = len(steps)
    if n == 0:
        return
    sw, _ = _slide_dims(prs)
    left = int(Inches(_MARGIN_IN))
    content_w = int(sw - Inches(_MARGIN_IN * 2))
    area_top = int(Inches(top_in))
    area_h = int(Inches(_AREA_BOTTOM_IN)) - area_top
    row_h = area_h // n
    badge = int(Inches(0.6))
    for idx, st in enumerate(steps):
        y = area_top + row_h * idx
        cy = y + (row_h - badge) // 2
        # 번호 배지.
        _add_rect(slide, left, cy, badge, badge, theme.accent)
        _add_text(slide, left, cy, badge, badge, str(idx + 1), size=22,
                  color=theme.white, font=theme.heading_font, bold=True,
                  align="center", anchor="middle")
        text_x = left + badge + int(Inches(0.35))
        text_w = content_w - badge - int(Inches(0.35))
        _add_text(slide, text_x, y + int(Inches(0.08)), text_w, Inches(0.45),
                  st.get("title", ""), size=17, color=theme.primary,
                  font=theme.heading_font, bold=True, align="left", anchor="middle")
        detail = st.get("detail")
        if detail:
            _add_text(slide, text_x, y + int(Inches(0.55)), text_w,
                      int(row_h) - int(Inches(0.6)), detail, size=13,
                      color=theme.muted, font=theme.body_font, align="left",
                      spacing=1.1)
        # 행 사이 헤어라인.
        if idx < n - 1:
            from pptx.util import Pt
            _add_rect(slide, text_x, y + row_h - int(Pt(0.5)), text_w, int(Pt(0.5)),
                      theme.hairline)


def _render_quote_slide(prs, theme: Theme, data: dict, page_no: int) -> None:
    """인용/인사이트 콜아웃 — 핵심 한 줄을 풀 슬라이드로 강조."""
    from pptx.util import Inches, Pt

    slide = prs.slides.add_slide(_blank_layout(prs))
    sw, sh = _slide_dims(prs)
    m = int(Inches(_MARGIN_IN))
    # 중성 배경 + 좌측 액센트 룰.
    _add_rect(slide, 0, 0, sw, sh, theme.surface)
    _add_rect(slide, 0, 0, int(Inches(0.16)), sh, theme.accent)
    # 거대한 옅은 인용부호.
    _add_text(slide, m, int(Inches(0.7)), int(Inches(3.0)), int(Inches(2.2)), "“",
              size=200, color=theme.accent_soft, font=theme.heading_font, bold=True,
              align="left", anchor="top", wrap=False)
    # 핵심 문장(중앙).
    _add_text(slide, m + int(Inches(0.3)), int(Inches(2.5)),
              sw - int(Inches(_MARGIN_IN * 2)) - int(Inches(0.3)), int(Inches(2.6)),
              data.get("text", ""), size=30, color=theme.primary,
              font=theme.heading_font, bold=True, align="left", anchor="middle",
              spacing=1.15)
    attribution = data.get("attribution")
    if attribution:
        _add_rect(slide, m + int(Inches(0.32)), int(Inches(5.4)), int(Inches(0.7)),
                  int(Pt(3)), theme.accent)
        _add_text(slide, m + int(Inches(0.32)), int(Inches(5.55)), int(Inches(8.0)),
                  int(Inches(0.5)), f"— {attribution}", size=15, color=theme.muted,
                  font=theme.body_font, align="left")
    _add_footer(slide, prs, theme, page_no)
