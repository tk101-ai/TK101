"""사용자 업로드 자료 텍스트 추출 (FR-03).

지원 포맷: PDF, DOCX, XLSX, CSV, TXT, PPTX
- T2 (`backend/app/services/nas_search/text_extractor.py`) 의 PDF/DOCX/PPTX 추출기 재사용
- XLSX, CSV, TXT 는 본 모듈에서 추가
- chunk_text 도 T2 의 청킹 정책 재사용 (1500자 윈도우, 150자 overlap)
"""
from __future__ import annotations

import csv
import io
import logging
from dataclasses import dataclass
from pathlib import Path

from app.services.nas_search.text_extractor import chunk_text as _t2_chunk_text
from app.services.nas_search.text_extractor import extract_text as _t2_extract_path

logger = logging.getLogger(__name__)

SUPPORTED_EXTENSIONS = {".pdf", ".docx", ".xlsx", ".csv", ".txt", ".pptx"}
DEFAULT_CHUNK_CHARS = 1500
DEFAULT_OVERLAP_CHARS = 150


@dataclass(frozen=True)
class ExtractedDocument:
    """업로드 자료 1개의 추출 결과."""

    filename: str
    extension: str
    text: str
    chunks: list[str]


def is_supported(filename: str) -> bool:
    return Path(filename).suffix.lower() in SUPPORTED_EXTENSIONS


def extract_from_path(file_path: str) -> str:
    """T2 추출기 위임. PDF/DOCX/PPTX는 T2 모듈, XLSX/CSV/TXT는 본 모듈."""
    ext = Path(file_path).suffix.lower()
    if ext in {".pdf", ".docx", ".pptx"}:
        return _t2_extract_path(file_path)
    if ext == ".xlsx":
        return _extract_xlsx_path(file_path)
    if ext == ".csv":
        return _extract_csv_path(file_path)
    if ext == ".txt":
        return _extract_txt_path(file_path)
    logger.debug("지원하지 않는 확장자: %s", file_path)
    return ""


def extract_from_bytes(file_bytes: bytes, filename: str) -> str:
    """multipart upload 처럼 메모리에서 직접 추출.

    PDF/DOCX/PPTX 는 BytesIO 래핑, XLSX는 openpyxl, CSV/TXT는 디코딩만.
    """
    ext = Path(filename).suffix.lower()
    try:
        if ext == ".pdf":
            return _extract_pdf_bytes(file_bytes)
        if ext == ".docx":
            return _extract_docx_bytes(file_bytes)
        if ext == ".pptx":
            return _extract_pptx_bytes(file_bytes)
        if ext == ".xlsx":
            return _extract_xlsx_bytes(file_bytes)
        if ext == ".csv":
            return _extract_csv_bytes(file_bytes)
        if ext == ".txt":
            return _extract_txt_bytes(file_bytes)
    except Exception as exc:  # noqa: BLE001 - 외부 라이브러리 다양한 예외 흡수
        logger.warning("텍스트 추출 실패: %s (%s)", filename, exc)
        return ""
    logger.debug("지원하지 않는 확장자: %s", filename)
    return ""


# --- PDF / DOCX / PPTX (bytes) ---------------------------------------------


def _extract_pdf_bytes(file_bytes: bytes) -> str:
    from pdfminer.high_level import extract_text as pdf_extract

    text = pdf_extract(io.BytesIO(file_bytes)) or ""
    return text.replace("\x00", "").strip()


def _extract_docx_bytes(file_bytes: bytes) -> str:
    from docx import Document

    doc = Document(io.BytesIO(file_bytes))
    parts: list[str] = []
    for para in doc.paragraphs:
        if para.text:
            parts.append(para.text)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                if cell.text:
                    parts.append(cell.text)
    return "\n".join(parts).replace("\x00", "").strip()


def _extract_pptx_bytes(file_bytes: bytes) -> str:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            text = getattr(shape, "text", None)
            if text:
                parts.append(text)
    return "\n".join(parts).replace("\x00", "").strip()


# --- XLSX -------------------------------------------------------------------


def _extract_xlsx_path(path: str) -> str:
    with open(path, "rb") as f:
        return _extract_xlsx_bytes(f.read())


def _extract_xlsx_bytes(file_bytes: bytes) -> str:
    """모든 시트의 셀 값을 행 단위 탭 구분 텍스트로 직렬화."""
    from openpyxl import load_workbook

    wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    parts: list[str] = []
    for sheet in wb.worksheets:
        parts.append(f"[시트: {sheet.title}]")
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(c.strip() for c in cells):
                parts.append("\t".join(cells))
    return "\n".join(parts).replace("\x00", "").strip()


# --- CSV --------------------------------------------------------------------


def _extract_csv_path(path: str) -> str:
    with open(path, "rb") as f:
        return _extract_csv_bytes(f.read())


def _extract_csv_bytes(file_bytes: bytes) -> str:
    text = _decode_text(file_bytes)
    if not text:
        return ""
    reader = csv.reader(io.StringIO(text))
    parts: list[str] = []
    for row in reader:
        if any(c.strip() for c in row):
            parts.append("\t".join(row))
    return "\n".join(parts).strip()


# --- TXT --------------------------------------------------------------------


def _extract_txt_path(path: str) -> str:
    with open(path, "rb") as f:
        return _extract_txt_bytes(f.read())


def _extract_txt_bytes(file_bytes: bytes) -> str:
    return _decode_text(file_bytes).replace("\x00", "").strip()


def _decode_text(file_bytes: bytes) -> str:
    """UTF-8 우선, 실패 시 CP949(한글 Windows) 폴백."""
    for encoding in ("utf-8", "utf-8-sig", "cp949", "euc-kr", "latin-1"):
        try:
            return file_bytes.decode(encoding)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


# --- 청킹 -------------------------------------------------------------------


def chunk_text(text: str) -> list[str]:
    """T2 청킹 정책 재사용. 1500자 윈도우 + 150자 overlap."""
    return _t2_chunk_text(
        text,
        chunk_chars=DEFAULT_CHUNK_CHARS,
        overlap_chars=DEFAULT_OVERLAP_CHARS,
    )


def extract_and_chunk(
    file_bytes: bytes, filename: str
) -> ExtractedDocument:
    """업로드 1건 → 추출 텍스트 + 청크. 라우터에서 form_data_sources 저장에 사용."""
    text = extract_from_bytes(file_bytes, filename)
    chunks = chunk_text(text) if text else []
    return ExtractedDocument(
        filename=filename,
        extension=Path(filename).suffix.lower(),
        text=text,
        chunks=chunks,
    )
