"""운영보고서 PPTX 양식에 표·요약을 채워 새 PPTX bytes 생성.

양식의 디자인/레이아웃을 그대로 두고, 표 셀과 요약 셀의 **텍스트만** 채운다
(서식=기존 run의 폰트 유지). 데이터 행 수가 양식의 빈 행 수와 다르면 마지막
행을 복제/삭제해 정확히 맞춘다.

슬라이드 식별은 인덱스가 아니라 **본문 텍스트 + 표 헤더**로 분류해, 양식이 약간
바뀌어도(슬라이드 추가 등) 견고하게 동작한다. 1차 범위: 운영 상세 표(중화권
OTA/따종디엔핑, 북미) + 운영 요약(진행제품/배포수량/배포데이터 자동, 홍보방향/
이슈는 AI 초안 주입). Top3·댓글 분석 슬라이드는 손대지 않는다(수동 유지).
"""
from __future__ import annotations

import io
import re
from copy import deepcopy

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.util import Emu, Pt

from .report_builder import (
    CHINA_DZDP_PLATFORMS,
    CHINA_OTA_PLATFORMS,
    china_detail_rows,
    na_detail_rows,
)
from .sheet_parser import DistRecord, MonthSection


# --- 셀 텍스트 쓰기(서식 보존) ----------------------------------------------
def _set_cell_text(cell, text: str) -> None:
    """첫 단락·첫 run의 폰트를 유지하며 텍스트 교체. 여분 run/단락은 제거.
    개행(\n)이 있으면 단락을 추가해 줄을 나눈다."""
    tf = cell.text_frame
    lines = (text or "").split("\n")
    p0 = tf.paragraphs[0]
    # 첫 단락: 첫 run 유지, 나머지 run 제거
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = lines[0]
    # 여분 단락 제거
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    # 추가 줄: 첫 단락을 복제해 폰트 유지
    for line in lines[1:]:
        new_p = deepcopy(p0._p)
        p0._p.getparent().append(new_p)
        # 복제 단락의 첫 run 텍스트만 교체
        from pptx.text.text import _Paragraph
        para = _Paragraph(new_p, p0._parent)
        if para.runs:
            para.runs[0].text = line
            for r in para.runs[1:]:
                r._r.getparent().remove(r._r)


def _set_cell_link(cell, label: str, url: str) -> None:
    """셀에 하이퍼링크 텍스트(예: '바로가기') 설정."""
    tf = cell.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        run = p0.runs[0]
        run.text = label
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        run = p0.add_run()
        run.text = label
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    if url:
        try:
            run.hyperlink.address = url
        except Exception:  # noqa: BLE001 — 링크 실패해도 라벨은 남김
            pass


# --- 표 행 수 맞추기 ---------------------------------------------------------
def _ensure_data_rows(table, n_data: int) -> None:
    """헤더(0행) 아래 데이터 행 수를 정확히 n_data 로 맞춘다(마지막 행 복제/삭제)."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    cur = len(trs) - 1  # 헤더 제외
    if cur < 1:
        return
    if cur < n_data:
        template_tr = trs[-1]
        for _ in range(n_data - cur):
            tbl.append(deepcopy(template_tr))
    elif cur > n_data:
        for tr in trs[1 + n_data:]:
            tbl.remove(tr)


def _fill_detail_table(table, rows: list[list]) -> None:
    """데이터 행을 채운다. 양식의 빈 데이터 셀은 run 폰트가 없어 채우면 기본(큰)
    폰트로 렌더돼 행이 부풀고 표가 슬라이드를 넘친다 → 데이터 셀에 명시적 폰트
    (검정 10.5pt 나눔고딕, 실제 보고서와 동급)를 줘 압축한다. (헤더는 흰 글씨라
    복사하면 데이터가 흰색으로 묻혀 안 됨 → 명시 폰트 사용.)"""
    _ensure_data_rows(table, len(rows))
    ncols = len(table.columns)
    for ri, rowvals in enumerate(rows, start=1):
        for ci, val in enumerate(rowvals):
            if ci >= ncols:
                break
            cell = table.cell(ri, ci)
            is_link = isinstance(val, tuple)
            if is_link:
                _set_cell_link(cell, val[0], val[1])
            else:
                _set_cell_text(cell, str(val))
            _apply_data_font(cell, link=is_link)


# 데이터 셀 기본 폰트(실제 보고서 데이터 셀 ~10.5~11pt 나눔고딕에 맞춤).
_DATA_FONT_PT = 10.5
_DATA_FONT_NAME = "나눔고딕"


def _apply_data_font(cell, *, link: bool = False) -> None:
    """데이터 셀 run 에 명시적 폰트(크기/이름/검정)를 적용해 행 부풀음 방지.
    링크 셀은 색을 건드리지 않아 하이퍼링크 스타일(파랑)을 보존한다."""
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            run.font.size = Pt(_DATA_FONT_PT)
            run.font.name = _DATA_FONT_NAME
            run.font.bold = False
            if not link:
                run.font.color.rgb = RGBColor(0x00, 0x00, 0x00)


def _copy_run_format(donor_cell, target_cell, *, unbold: bool = True) -> None:
    """donor 셀 첫 run의 폰트(rPr)를 target 셀 모든 run에 복사. 양식의 빈 값 셀은
    run 서식이 없어 기본(큰) 폰트로 렌더돼 짧은 행에서 잘리므로, 라벨 셀의 서식을
    빌려 일관된 크기로 채운다. 값은 라벨처럼 굵게 두지 않도록 bold는 끈다."""
    druns = donor_cell.text_frame.paragraphs[0].runs
    if not druns:
        return
    drpr = druns[0]._r.find(qn("a:rPr"))
    if drpr is None:
        return
    for para in target_cell.text_frame.paragraphs:
        for run in para.runs:
            old = run._r.find(qn("a:rPr"))
            if old is not None:
                run._r.remove(old)
            new = deepcopy(drpr)
            if unbold and new.get("b") is not None:
                del new.attrib["b"]
            run._r.insert(0, new)


def _enable_shrink_to_fit(cell) -> None:
    """셀 텍스트가 길면 칸에 맞게 자동 축소(normAutofit). AI 서술이 길어도 칸 밖으로
    삐져나오지 않게 한다(파워포인트/리브레오피스가 열 때 폰트 스케일 계산)."""
    bodyPr = cell.text_frame._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit", "a:noAutofit"):
        e = bodyPr.find(qn(tag))
        if e is not None:
            bodyPr.remove(e)
    bodyPr.append(bodyPr.makeelement(qn("a:normAutofit"), {}))


def _set_no_bullet(para) -> None:
    """단락의 자동번호/불릿 제거(buNone). 템플릿 자동번호와 텍스트 번호가 겹쳐
    '1. 1.'로 중복되던 문제 방지."""
    pPr = para._p.get_or_add_pPr()
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        e = pPr.find(qn(tag))
        if e is not None:
            pPr.remove(e)
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def _fill_promo_cell(cell, directions: list) -> None:
    """홍보 방향 셀을 기존 보고서 형식으로 렌더: 파란 굵은 'N. 제목' + 검정 '• 상세'
    + 방향 사이 빈 줄. 자동번호는 끄고 번호를 직접 넣어 중복을 막는다."""
    tf = cell.text_frame
    tf.clear()
    first = True
    n = len(directions)
    for i, d in enumerate(directions, 1):
        head = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        _set_no_bullet(head)
        r = head.add_run()
        r.text = f"{i}. {d.get('title', '').strip()}"
        r.font.bold = True
        r.font.size = Pt(12)
        r.font.color.theme_color = MSO_THEME_COLOR.ACCENT_1  # 기존 보고서 파란색
        for det in d.get("details", []):
            p = tf.add_paragraph()
            _set_no_bullet(p)
            rd = p.add_run()
            rd.text = f"• {str(det).strip()}"
            rd.font.bold = False
            rd.font.size = Pt(12)
            rd.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
        if i < n:  # 방향 사이 빈 줄
            _set_no_bullet(tf.add_paragraph())
    _enable_shrink_to_fit(cell)


def _fill_narrative(table, narrative: dict) -> None:
    """요약표의 홍보 방향(구조화)·진행 이슈사항(한 문장)을 AI 초안으로 채움."""
    if not narrative:
        return
    directions = narrative.get("directions")
    issue = narrative.get("issue")
    for ri in range(len(table.rows)):
        label = re.sub(r"\s+", "", table.cell(ri, 0).text)
        if label == "홍보방향" and directions:
            _fill_promo_cell(table.cell(ri, 1), directions)
        elif label == "진행이슈사항" and issue:
            cell = table.cell(ri, 1)
            _set_cell_text(cell, f"[초안] {issue}")
            _copy_run_format(table.cell(ri, 0), cell)
            _enable_shrink_to_fit(cell)


def _fill_summary_table(table, values: dict) -> None:
    """5x2 요약표: col0 라벨과 매칭되는 col1 셀을 채움(라벨 정규화로 줄바꿈 무시).
    값 셀은 자동축소를 켜 긴 AI 서술도 칸 안에 맞춘다."""
    for ri in range(len(table.rows)):
        label = re.sub(r"\s+", "", table.cell(ri, 0).text)
        for key, val in values.items():
            if re.sub(r"\s+", "", key) == label and val:
                cell = table.cell(ri, 1)
                _set_cell_text(cell, val)
                _copy_run_format(table.cell(ri, 0), cell)
                _enable_shrink_to_fit(cell)
                break


# --- 슬라이드 분류 -----------------------------------------------------------
def _slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text)
        except Exception:  # noqa: BLE001
            pass
    return "\n".join(parts)


def _tables(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


def _is_summary(table) -> bool:
    return len(table.columns) == 2 and "진행" in table.cell(0, 0).text and "제품" in table.cell(0, 0).text


def _is_detail(table) -> bool:
    hdr = "".join(table.cell(0, c).text for c in range(len(table.columns)))
    return "계정명" in hdr and "팔로워" in hdr


# --- 표지/기준일 ------------------------------------------------------------
# Top3·댓글 슬라이드(이전월 캡처/댓글 잔존) 식별 키워드.
_CONTENT_ANALYSIS_KEYS = ("우수 콘텐츠", "댓글 분석")
# 운영 리뷰 슬라이드(AS-IS/TO-BE 구조 보존) 식별 키워드.
_REVIEW_KEYS = ("운영 리뷰", "운영 제안", "AS-IS", "TO-BE")
# 제목 영역 아래(분석 박스 시작). 이 위(top 작은 값)는 헤더/섹션 제목이라 보존.
_TITLE_BAND_IN = 1.3
# 콘텐츠 영역 시작 높이(in) — 이 아래가 캡처·댓글 등 잔존 미디어/텍스트.
_CONTENT_TOP_IN = 2.8
_EMU_PER_IN = 914400


def _clear_stale_media(prs) -> None:
    """Top3·댓글·리뷰 슬라이드의 **이전월 스크린샷·댓글** 잔존물을 제거(자동 비우기).

    양식이 5월 보고서 기반이라 이 슬라이드들엔 5월 캡처/댓글이 박혀 있어 다른 달
    보고서에 틀린 내용이 보인다. 헤더·제목·빈 분석박스(상단)는 보존하고 콘텐츠
    영역(분석박스 아래)의 잔존 도형을 지워 마케터가 해당 월 자료를 채울 빈 영역으로
    남긴다.
    - Top3·댓글: 콘텐츠 영역의 모든 도형(사진·그룹·말풍선·댓글 텍스트) 제거.
    - 운영 리뷰: AS-IS/TO-BE 구조는 보존하고 사진/그룹(이미지)만 제거.
    """
    title_band = Emu(int(_TITLE_BAND_IN * _EMU_PER_IN))
    content_top = Emu(int(_CONTENT_TOP_IN * _EMU_PER_IN))
    for slide in prs.slides:
        text = _slide_text(slide)
        is_content = any(k in text for k in _CONTENT_ANALYSIS_KEYS)
        is_review = any(k in text for k in _REVIEW_KEYS)
        if not (is_content or is_review):
            continue
        for sh in list(slide.shapes):
            if sh.top is None or sh.top < title_band:
                continue  # 헤더·섹션 제목 보존
            is_media = sh.shape_type in (MSO_SHAPE_TYPE.PICTURE, MSO_SHAPE_TYPE.GROUP)
            if sh.top >= content_top:
                # 콘텐츠 영역: 캡처·말풍선·댓글 잔존물.
                # 리뷰는 이미지만 제거(AS-IS/TO-BE 구조 보존), 그 외는 전부 제거.
                if is_review and not is_content and not is_media:
                    continue
                sh._element.getparent().remove(sh._element)
            elif is_content:
                # 분석 박스 영역(제목~콘텐츠 사이): 이미지는 제거, 텍스트 박스는
                # 비운다(이전월 분석 텍스트 잔존 방지, 박스 자체는 빈 채로 보존).
                if is_media:
                    sh._element.getparent().remove(sh._element)
                elif getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip():
                    sh.text_frame.clear()


def _update_cover(prs, month: int, basis_date: str | None) -> None:
    """표지·부제목·각주의 '월'/'기준일' 텍스트를 보고 월로 갱신.

    양식이 5월 보고서 기반이라 '5월 운영 내용 요약', '5월', '운영보고서' 등 월 라벨이
    텍스트 상자(표 셀 아님)에 박혀 있다. 표 셀은 GraphicFrame 이라 여기서 안 건드린다
    (데이터 날짜 오염 없음). 월 표기가 있는 제목/부제목 run 의 'N월'만 교체한다.
    """
    date_re = r"\d{4}\s*년\s*\d+\s*월\s*\d+\s*일"
    for slide in prs.slides:
        for sh in slide.shapes:
            if not (getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()):
                continue
            for para in sh.text_frame.paragraphs:
                # 월 숫자와 '월'이 별도 run 으로 쪼개질 수 있어 **단락 단위**로 합쳐 치환.
                ptext = "".join(r.text for r in para.runs)
                if not ptext.strip():
                    continue
                if basis_date and re.search(date_re, ptext):
                    new = re.sub(date_re, basis_date, ptext)
                elif "기준으로 기재" in ptext:
                    continue  # 기준일 노트는 basis_date 없으면 건드리지 않음
                elif re.search(r"\d+\s*월", ptext):
                    new = re.sub(r"\d+\s*월", f"{month}월", ptext)
                else:
                    continue
                if new != ptext and para.runs:
                    # 첫 run 에 합친 텍스트를 넣고 나머지 run 제거(제목/부제목은 균일 서식).
                    para.runs[0].text = new
                    for r in para.runs[1:]:
                        r._r.getparent().remove(r._r)


# --- 메인 -------------------------------------------------------------------
def fill_report(
    *,
    template_bytes: bytes,
    month: int,
    china_records: list[DistRecord],
    na_records: list[DistRecord],
    china_summary_vals: dict,
    na_summary_vals: dict,
    china_narrative: dict | None = None,
    na_narrative: dict | None = None,
    basis_date: str | None = None,
) -> bytes:
    """양식 bytes + 채움 데이터 → 채워진 PPTX bytes.

    슬라이드를 순회하며 직전 섹션 구분('1. 중화권'/'2. 북미')으로 region을 추적하고,
    표 종류(요약/상세)와 본문 텍스트(따종디엔핑 여부)로 무엇을 채울지 결정한다.
    """
    prs = Presentation(io.BytesIO(template_bytes))
    _update_cover(prs, month, basis_date)
    _clear_stale_media(prs)  # 이전월 캡처(Top3·댓글·리뷰) 자동 제거

    region = None  # "china" | "na" | "review"
    for slide in prs.slides:
        text = _slide_text(slide)
        if re.search(r"1\s*[.．]?\s*중화권", text):
            region = "china"
        elif re.search(r"2\s*[.．]?\s*북미", text):
            region = "na"
        elif re.search(r"3\s*[.．]?\s*리뷰", text):
            region = "review"

        for table in _tables(slide):
            if _is_summary(table):
                is_na = region == "na"
                # 계산 값(진행제품/배포수량/배포데이터)은 결정적으로 채우고,
                # 홍보 방향/진행 이슈는 AI 초안을 기존 보고서 형식으로 렌더.
                _fill_summary_table(table, na_summary_vals if is_na else china_summary_vals)
                _fill_narrative(table, (na_narrative if is_na else china_narrative) or {})
            elif _is_detail(table):
                if region == "na":
                    _fill_detail_table(table, na_detail_rows(na_records))
                elif region == "china":
                    if "따종" in text:
                        _fill_detail_table(
                            table, china_detail_rows(china_records, CHINA_DZDP_PLATFORMS)
                        )
                    else:
                        _fill_detail_table(
                            table, china_detail_rows(china_records, CHINA_OTA_PLATFORMS)
                        )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
