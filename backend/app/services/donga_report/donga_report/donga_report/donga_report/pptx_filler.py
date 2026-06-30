"""운영보고서 PPTX 양식에 표·요약을 채워 새 PPTX bytes 생성.

양식의 디자인/레이아웃을 그대로 두고, 표 셀과 요약 셀의 **텍스트만** 채운다
(서식=기존 run의 폰트 유지). 데이터 행 수가 양식의 빈 행 수와 다르면 마지막
행을 복제/삭제해 정확히 맞춘다.

슬라이드 식별은 인덱스가 아니라 **본문 텍스트 + 표 헤더**로 분류해, 양식이 약간
바뀌어도(슬라이드 추가 등) 견고하게 동작한다. 1차 범위: 운영 상세 표(중화권
OTA/따종디엔핑, 북미) + 운영 요약(진행제품/배포수량/배포데이터 자동, 홍보방향/
이슈는 AI 초안 주입). Top3·댓글 분석 슬라이드는 손대지 않는다(수동 유지).
"""
from __future__ import annotations

import io
import re
from copy import deepcopy

from pptx import Presentation
from pptx.oxml.ns import qn

from .report_builder import (
    CHINA_DZDP_PLATFORMS,
    CHINA_OTA_PLATFORMS,
    china_detail_rows,
    na_detail_rows,
)
from .sheet_parser import DistRecord, MonthSection


# --- 셀 텍스트 쓰기(서식 보존) ----------------------------------------------
def _set_cell_text(cell, text: str) -> None:
    """첫 단락·첫 run의 폰트를 유지하며 텍스트 교체. 여분 run/단락은 제거.
    개행(\n)이 있으면 단락을 추가해 줄을 나눈다."""
    tf = cell.text_frame
    lines = (text or "").split("\n")
    p0 = tf.paragraphs[0]
    # 첫 단락: 첫 run 유지, 나머지 run 제거
    if p0.runs:
        p0.runs[0].text = lines[0]
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        p0.add_run().text = lines[0]
    # 여분 단락 제거
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    # 추가 줄: 첫 단락을 복제해 폰트 유지
    for line in lines[1:]:
        new_p = deepcopy(p0._p)
        p0._p.getparent().append(new_p)
        # 복제 단락의 첫 run 텍스트만 교체
        from pptx.text.text import _Paragraph
        para = _Paragraph(new_p, p0._parent)
        if para.runs:
            para.runs[0].text = line
            for r in para.runs[1:]:
                r._r.getparent().remove(r._r)


def _set_cell_link(cell, label: str, url: str) -> None:
    """셀에 하이퍼링크 텍스트(예: '바로가기') 설정."""
    tf = cell.text_frame
    p0 = tf.paragraphs[0]
    if p0.runs:
        run = p0.runs[0]
        run.text = label
        for r in p0.runs[1:]:
            r._r.getparent().remove(r._r)
    else:
        run = p0.add_run()
        run.text = label
    for extra in tf.paragraphs[1:]:
        extra._p.getparent().remove(extra._p)
    if url:
        try:
            run.hyperlink.address = url
        except Exception:  # noqa: BLE001 — 링크 실패해도 라벨은 남김
            pass


# --- 표 행 수 맞추기 ---------------------------------------------------------
def _ensure_data_rows(table, n_data: int) -> None:
    """헤더(0행) 아래 데이터 행 수를 정확히 n_data 로 맞춘다(마지막 행 복제/삭제)."""
    tbl = table._tbl
    trs = tbl.findall(qn("a:tr"))
    cur = len(trs) - 1  # 헤더 제외
    if cur < 1:
        return
    if cur < n_data:
        template_tr = trs[-1]
        for _ in range(n_data - cur):
            tbl.append(deepcopy(template_tr))
    elif cur > n_data:
        for tr in trs[1 + n_data:]:
            tbl.remove(tr)


def _fill_detail_table(table, rows: list[list]) -> None:
    _ensure_data_rows(table, len(rows))
    ncols = len(table.columns)
    for ri, rowvals in enumerate(rows, start=1):
        for ci, val in enumerate(rowvals):
            if ci >= ncols:
                break
            cell = table.cell(ri, ci)
            if isinstance(val, tuple):
                _set_cell_link(cell, val[0], val[1])
            else:
                _set_cell_text(cell, str(val))


def _fill_summary_table(table, values: dict) -> None:
    """5x2 요약표: col0 라벨과 매칭되는 col1 셀을 채움(라벨 정규화로 줄바꿈 무시)."""
    for ri in range(len(table.rows)):
        label = re.sub(r"\s+", "", table.cell(ri, 0).text)
        for key, val in values.items():
            if re.sub(r"\s+", "", key) == label and val:
                _set_cell_text(table.cell(ri, 1), val)
                break


# --- 슬라이드 분류 -----------------------------------------------------------
def _slide_text(slide) -> str:
    parts = []
    for sh in slide.shapes:
        try:
            if sh.has_text_frame and sh.text_frame.text.strip():
                parts.append(sh.text_frame.text)
        except Exception:  # noqa: BLE001
            pass
    return "\n".join(parts)


def _tables(slide):
    return [sh.table for sh in slide.shapes if sh.has_table]


def _is_summary(table) -> bool:
    return len(table.columns) == 2 and "진행" in table.cell(0, 0).text and "제품" in table.cell(0, 0).text


def _is_detail(table) -> bool:
    hdr = "".join(table.cell(0, c).text for c in range(len(table.columns)))
    return "계정명" in hdr and "팔로워" in hdr


# --- 표지/기준일 ------------------------------------------------------------
def _update_cover(prs, month: int, basis_date: str | None) -> None:
    """표지·각주의 '월'/'기준일' 텍스트 갱신(있을 때만)."""
    for slide in prs.slides:
        for sh in slide.shapes:
            if not (getattr(sh, "has_text_frame", False) and sh.text_frame.text.strip()):
                continue
            for para in sh.text_frame.paragraphs:
                for run in para.runs:
                    t = run.text
                    if "운영보고서" in t:
                        run.text = re.sub(r"\d+\s*월", f"{month}월", t)
                    if basis_date and "기준으로 기재" in t:
                        run.text = re.sub(
                            r"\d{4}년\s*\d+월\s*\d+일", basis_date, t
                        )


# --- 메인 -------------------------------------------------------------------
def fill_report(
    *,
    template_bytes: bytes,
    month: int,
    china_records: list[DistRecord],
    na_records: list[DistRecord],
    china_summary_vals: dict,
    na_summary_vals: dict,
    china_narrative: dict | None = None,
    na_narrative: dict | None = None,
    basis_date: str | None = None,
) -> bytes:
    """양식 bytes + 채움 데이터 → 채워진 PPTX bytes.

    슬라이드를 순회하며 직전 섹션 구분('1. 중화권'/'2. 북미')으로 region을 추적하고,
    표 종류(요약/상세)와 본문 텍스트(따종디엔핑 여부)로 무엇을 채울지 결정한다.
    """
    prs = Presentation(io.BytesIO(template_bytes))
    _update_cover(prs, month, basis_date)

    china_sum = {**china_summary_vals, **(china_narrative or {})}
    na_sum = {**na_summary_vals, **(na_narrative or {})}

    region = None  # "china" | "na" | "review"
    for slide in prs.slides:
        text = _slide_text(slide)
        if re.search(r"1\s*[.．]?\s*중화권", text):
            region = "china"
        elif re.search(r"2\s*[.．]?\s*북미", text):
            region = "na"
        elif re.search(r"3\s*[.．]?\s*리뷰", text):
            region = "review"

        for table in _tables(slide):
            if _is_summary(table):
                _fill_summary_table(table, na_sum if region == "na" else china_sum)
            elif _is_detail(table):
                if region == "na":
                    _fill_detail_table(table, na_detail_rows(na_records))
                elif region == "china":
                    if "따종" in text:
                        _fill_detail_table(
                            table, china_detail_rows(china_records, CHINA_DZDP_PLATFORMS)
                        )
                    else:
                        _fill_detail_table(
                            table, china_detail_rows(china_records, CHINA_OTA_PLATFORMS)
                        )

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()
