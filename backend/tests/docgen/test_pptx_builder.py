"""pptx_builder 고품질 렌더 — 회귀/견고성 테스트.

실제 .pptx 바이트를 만들어 python-pptx 로 다시 열어 슬라이드 구성/크래시無를
검증한다(시각 품질 자체는 라이브 산출물로 별도 확인). 핵심은 **자동 생성물이
어떤 입력에도 통째로 깨지지 않는다**는 보장.
"""
from __future__ import annotations

import io

from pptx import Presentation

from app.services.docgen.pptx_builder import build_pptx
from app.services.docgen.theme import get_theme

SECTIONS = [
    {
        "heading": "시장 개요",
        "body": "국내 시장은 성장세다.\n\n- 핵심 포인트 1\n  - 하위 근거\n- 핵심 포인트 2",
    },
    {
        "heading": "실적 표",
        "body": "| 분기 | 매출 | 성장률 |\n|---|---|---|\n| 1Q | 10 | 5% |\n| 2Q | 14 | 40% |",
    },
    {
        "heading": "추세 차트",
        "body": '```chart\n{"type":"column","title":"분기 매출",'
        '"categories":["1Q","2Q"],"series":[{"name":"매출","values":[10,14]}]}\n```',
    },
]


def _open(data: bytes) -> Presentation:
    assert isinstance(data, bytes) and len(data) > 1000
    return Presentation(io.BytesIO(data))


def test_builds_valid_pptx_with_all_block_types() -> None:
    prs = _open(build_pptx("2026 사업 제안서", SECTIONS))
    # 표지 + 목차 + (섹션 구분 + 본문) — 섹션 3개라 충분한 슬라이드 수.
    assert len(prs.slides) >= 1 + 1 + 3 * 2


def test_single_section_skips_agenda_and_dividers() -> None:
    prs = _open(build_pptx("단일", [SECTIONS[0]]))
    # 표지 + 본문 1장(목차/구분 없음).
    assert len(prs.slides) == 2


def test_empty_sections_still_renders_cover() -> None:
    prs = _open(build_pptx("제목만", []))
    assert len(prs.slides) == 1


def test_malformed_pie_chart_does_not_crash_deck() -> None:
    # 다중 시리즈 pie 는 python-pptx 가 거부 → 플레이스홀더 폴백, 덱은 정상.
    bad = [
        {
            "heading": "잘못된 차트",
            "body": '```chart\n{"type":"pie","title":"불량",'
            '"categories":["A","B"],"series":['
            '{"name":"s1","values":[1,2]},{"name":"s2","values":[3,4]}]}\n```',
        }
    ]
    prs = _open(build_pptx("폴백", bad))
    assert len(prs.slides) >= 2


def test_missing_heading_and_body_is_safe() -> None:
    prs = _open(build_pptx("", [{"heading": "", "body": ""}]))
    assert len(prs.slides) >= 1


def test_theme_exposes_derived_design_tokens() -> None:
    theme = get_theme()
    assert len(theme.series_palette) >= 4
    # 파생 토큰이 모두 RGB 3튜플.
    for tok in (theme.primary_deep, theme.accent_soft, theme.surface, theme.hairline):
        assert isinstance(tok, tuple) and len(tok) == 3
