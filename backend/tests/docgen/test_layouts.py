"""레이아웃 아카이브 — 파서 정규화 + pptx/docx 렌더 견고성 테스트.

핵심: ```layout``` 펜스가 정확히 파싱되고, 5종 아카이브가 유효한 .pptx 를 만들며,
불량 스펙은 조용히 건너뛰고, docx 는 데이터 유실 없이 폴백한다(어떤 입력도 안 깨짐).
"""
from __future__ import annotations

import io

from docx import Document
from pptx import Presentation

from app.services.docgen.docx_builder import build_docx
from app.services.docgen.markdown_blocks import parse_blocks
from app.services.docgen.pptx_builder import build_pptx


def _layout_body(json_str: str) -> str:
    return f"```layout\n{json_str}\n```"


KPI = '{"layout":"kpi","title":"핵심 지표","items":[{"value":"40%","label":"성장률","caption":"전년比"},{"value":"1,200","label":"신규 고객"}]}'
COMPARE = '{"layout":"compare","title":"비교","columns":[{"heading":"기존","points":["느림","비쌈"]},{"heading":"제안","points":["빠름","저렴"]}]}'
TIMELINE = '{"layout":"timeline","title":"로드맵","milestones":[{"label":"1단계","title":"착수","detail":"준비"},{"label":"2단계","title":"확장","detail":"이중화"}]}'
STEPS = '{"layout":"steps","title":"절차","steps":[{"title":"분석","detail":"현황 파악"},{"title":"설계","detail":"아키텍처"},{"title":"구현"}]}'
QUOTE = '{"layout":"quote","text":"속도가 곧 경쟁력이다","attribution":"전략기획팀"}'


# ---------------------------------------------------------------------------
# 파서 정규화
# ---------------------------------------------------------------------------
class TestParseLayout:
    def test_kpi_parsed(self) -> None:
        blocks = parse_blocks(_layout_body(KPI))
        assert len(blocks) == 1 and blocks[0].kind == "layout"
        assert blocks[0].data["layout"] == "kpi"
        assert len(blocks[0].data["items"]) == 2

    def test_compare_requires_two_columns(self) -> None:
        one = '{"layout":"compare","columns":[{"heading":"A","points":["x"]}]}'
        assert parse_blocks(_layout_body(one)) == []  # 컬럼 1개 → 무시.

    def test_unknown_layout_skipped(self) -> None:
        assert parse_blocks(_layout_body('{"layout":"bogus","items":[]}')) == []

    def test_invalid_json_skipped(self) -> None:
        assert parse_blocks(_layout_body("{not json")) == []

    def test_kpi_caps_items(self) -> None:
        many = ('{"layout":"kpi","items":['
                + ",".join('{"value":"%d","label":"L%d"}' % (i, i) for i in range(8))
                + "]}")
        blocks = parse_blocks(_layout_body(many))
        assert len(blocks[0].data["items"]) == 4  # cap=4.

    def test_quote_requires_text(self) -> None:
        assert parse_blocks(_layout_body('{"layout":"quote","attribution":"x"}')) == []


# ---------------------------------------------------------------------------
# pptx 렌더
# ---------------------------------------------------------------------------
def _open_pptx(data: bytes) -> Presentation:
    assert isinstance(data, bytes) and len(data) > 1000
    return Presentation(io.BytesIO(data))


class TestPptxLayoutRender:
    def test_all_archetypes_render(self) -> None:
        sections = [
            {"heading": "지표", "body": "요약 문장.\n\n" + _layout_body(KPI)},
            {"heading": "비교", "body": _layout_body(COMPARE)},
            {"heading": "로드맵", "body": _layout_body(TIMELINE)},
            {"heading": "절차", "body": _layout_body(STEPS)},
            {"heading": "인사이트", "body": _layout_body(QUOTE)},
        ]
        prs = _open_pptx(build_pptx("레이아웃 데모", sections))
        # 표지 + 목차 + 섹션 5개(구분+레이아웃) → 충분한 슬라이드.
        assert len(prs.slides) >= 1 + 1 + 5 * 2

    def test_layout_only_section_no_crash(self) -> None:
        prs = _open_pptx(build_pptx("단일", [{"heading": "지표", "body": _layout_body(KPI)}]))
        assert len(prs.slides) >= 2

    def test_caption_folds_into_layout_slide(self) -> None:
        # 레이아웃 위 짧은 리드 문장은 같은 슬라이드 캡션 → 빈 슬라이드 없음.
        body = "이것은 리드 문장이다.\n\n" + _layout_body(KPI)
        prs = _open_pptx(build_pptx("캡션", [{"heading": "지표", "body": body}]))
        # 표지 + 레이아웃 1장 = 2장(리드용 별도 슬라이드 없음).
        assert len(prs.slides) == 2


# ---------------------------------------------------------------------------
# docx 폴백(데이터 유실 방지)
# ---------------------------------------------------------------------------
class TestDocxFallback:
    def test_layout_blocks_do_not_crash_docx(self) -> None:
        sections = [
            {"heading": "지표", "body": _layout_body(KPI)},
            {"heading": "절차", "body": _layout_body(STEPS)},
            {"heading": "인용", "body": _layout_body(QUOTE)},
        ]
        data = build_docx("폴백", sections)
        doc = Document(io.BytesIO(data))
        text = "\n".join(p.text for p in doc.paragraphs)
        # 핵심 내용이 평문으로 보존됐는지(유실 방지).
        assert "속도가 곧 경쟁력이다" in text
        assert "분석" in text
